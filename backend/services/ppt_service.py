import os
import io
import re
import json
import asyncio
import logging
import uuid
from typing import List, Optional, Dict, Any
from pydantic import BaseModel, Field

from services.ai_provider import ai_provider
from services.groq_service import robust_json_parser
from services.error_service import format_ai_exception_detail
from services.ppt_history_service import ppt_history_service
from fastapi import HTTPException

# ReportLab imports for Landscape PDF
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable, KeepTogether,
    Image as RLImage
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.pdfgen import canvas
from services.pdf_service import (
    UNICODE_FONT_NAME, UNICODE_BOLD_FONT_NAME,
    clean_md_to_reportlab, strip_emojis_for_pdf
)

# PPTX imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger("ppt_service")

# --- Schemas ---

class SlideItem(BaseModel):
    slide_number: int
    layout: str = "title_bullets" # cover, title_bullets, two_column, stat_highlight, process_timeline, quote_insight, split_image_text, thank_you
    category: str = "Concept"
    title: str
    subtitle: Optional[str] = None
    bullets: List[str] = []
    left_column: Optional[Dict[str, Any]] = None
    right_column: Optional[Dict[str, Any]] = None
    metrics: Optional[List[Dict[str, Any]]] = None
    timeline_steps: Optional[List[Dict[str, Any]]] = None
    quote: Optional[Dict[str, Any]] = None
    image_keyword: str = "education study concept"
    image_url: Optional[str] = None
    image_caption: Optional[str] = None
    speaker_notes: str = ""

class GeneratePPTRequest(BaseModel):
    topic: str = Field(..., example="Quantum Computing and Superposition")
    target_audience: str = Field(default="Class 10-12 / High School")
    num_slides: int = Field(default=8, ge=3, le=20)
    tone: str = Field(default="Engaging & Visual")
    language: str = Field(default="English")
    theme: str = Field(default="modern_navy") # modern_navy, emerald_sage, sunset_coral, dark_cyber, royal_purple, slate_academic
    teacher_guidance: Optional[str] = Field(default="")
    user_email: Optional[str] = None
    user_id: Optional[str] = None
    presenter_name: Optional[str] = None

class PresentationData(BaseModel):
    id: Optional[str] = None
    title: str
    subtitle: str
    topic: str
    target_audience: str
    num_slides: int
    theme: str = "modern_navy"
    language: str = "English"
    teacher_guidance: Optional[str] = None
    presenter_name: Optional[str] = None
    slides: List[SlideItem]
    created_at: Optional[str] = None

class RefineSlideRequest(BaseModel):
    slide: SlideItem
    instruction: str
    topic: str
    target_audience: str = "Class 10-12"

# Curated high-res educational photo fallback library
THEME_PRESETS: Dict[str, Dict[str, Any]] = {
    "modern_navy": {
        "primary": "#1E3A8A", # Blue 900
        "secondary": "#3B82F6", # Blue 500
        "accent": "#0D9488", # Teal 600
        "bg": "#F8FAFC",
        "card_bg": "#FFFFFF",
        "text_primary": "#0F172A",
        "text_secondary": "#475569",
        "rgb_primary": (30, 58, 138),
        "rgb_accent": (13, 148, 136),
        "rgb_card": (255, 255, 255),
        "rgb_bg": (248, 250, 252)
    },
    "emerald_sage": {
        "primary": "#065F46", # Emerald 800
        "secondary": "#10B981", # Emerald 500
        "accent": "#F59E0B", # Amber 500
        "bg": "#F0FDF4",
        "card_bg": "#FFFFFF",
        "text_primary": "#064E3B",
        "text_secondary": "#334155",
        "rgb_primary": (6, 95, 70),
        "rgb_accent": (245, 158, 11),
        "rgb_card": (255, 255, 255),
        "rgb_bg": (240, 253, 244)
    },
    "sunset_coral": {
        "primary": "#991B1B", # Red 800
        "secondary": "#EA580C", # Orange 600
        "accent": "#D97706", # Amber 600
        "bg": "#FFF7ED",
        "card_bg": "#FFFFFF",
        "text_primary": "#7C2D12",
        "text_secondary": "#431407",
        "rgb_primary": (153, 27, 27),
        "rgb_accent": (234, 88, 12),
        "rgb_card": (255, 255, 255),
        "rgb_bg": (255, 247, 237)
    },
    "dark_cyber": {
        "primary": "#0F172A", # Slate 900
        "secondary": "#6366F1", # Indigo 500
        "accent": "#06B6D4", # Cyan 500
        "bg": "#0B0F19",
        "card_bg": "#1E293B",
        "text_primary": "#F8FAFC",
        "text_secondary": "#94A3B8",
        "rgb_primary": (99, 102, 241),
        "rgb_accent": (6, 182, 212),
        "rgb_card": (30, 41, 59),
        "rgb_bg": (11, 15, 25)
    },
    "royal_purple": {
        "primary": "#581C87", # Purple 900
        "secondary": "#8B5CF6", # Purple 500
        "accent": "#EC4899", # Pink 500
        "bg": "#FAF5FF",
        "card_bg": "#FFFFFF",
        "text_primary": "#3B0764",
        "text_secondary": "#475569",
        "rgb_primary": (88, 28, 135),
        "rgb_accent": (236, 72, 153),
        "rgb_card": (255, 255, 255),
        "rgb_bg": (250, 245, 255)
    },
    "slate_academic": {
        "primary": "#1E293B", # Slate 800
        "secondary": "#475569", # Slate 600
        "accent": "#2563EB", # Blue 600
        "bg": "#F1F5F9",
        "card_bg": "#FFFFFF",
        "text_primary": "#0F172A",
        "text_secondary": "#334155",
        "rgb_primary": (30, 41, 59),
        "rgb_accent": (37, 99, 235),
        "rgb_card": (255, 255, 255),
        "rgb_bg": (241, 245, 249)
    }
}

_IMAGE_CACHE: Dict[str, bytes] = {}
_REAL_IMAGE_URL_CACHE: Dict[str, str] = {}

def _download_image_bytes(url: Optional[str]) -> Optional[bytes]:
    """Downloads image bytes from URL or parses base64 data URLs with in-memory caching."""
    if not url:
        return None
    url_str = str(url).strip()
    
    # 1. Direct device-uploaded image support (base64 Data URL)
    if "base64," in url_str:
        try:
            import base64
            b64_part = url_str.split("base64,")[-1].strip()
            return base64.b64decode(b64_part)
        except Exception as b64_err:
            logger.warning(f"Failed decoding device base64 image data URL: {b64_err}")
            return None

    # 2. Remote HTTP/HTTPS URL
    if not url_str.startswith(("http://", "https://")):
        return None

    if url_str in _IMAGE_CACHE:
        return _IMAGE_CACHE[url_str]

    try:
        import httpx
        headers = {"User-Agent": "DEVGYA-Educational-App/1.0 (https://devgya.in; contact@devgya.in)"}
        with httpx.Client(timeout=6.0, headers=headers, follow_redirects=True) as client:
            resp = client.get(url_str)
            if resp.status_code == 200 and len(resp.content) > 500:
                _IMAGE_CACHE[url_str] = resp.content
                return resp.content
    except Exception as e:
        logger.warning(f"Failed downloading slide image from {url_str}: {e}")
    return None


async def _resolve_real_topic_image(topic: str, slide_title: str, keyword: str) -> str:
    """
    Dynamically finds real, authentic educational images matching the specific topic & slide concept.
    Uses Wikimedia Commons / Wikipedia API for authentic diagrams, maps, and photographs,
    falling back to Pollinations AI for photorealistic topic diagrams.
    """
    cache_key = f"{topic}_{slide_title}_{keyword}".lower().strip()
    if cache_key in _REAL_IMAGE_URL_CACHE:
        return _REAL_IMAGE_URL_CACHE[cache_key]

    queries = [
        f"{topic} {keyword}".strip(),
        f"{keyword}".strip(),
        f"{topic} {slide_title}".strip(),
        f"{slide_title}".strip()
    ]
    headers = {"User-Agent": "DEVGYA-Educational-App/1.0 (https://devgya.in; contact@devgya.in)"}

    try:
        import httpx
        import urllib.parse
        async with httpx.AsyncClient(timeout=4.0, headers=headers, follow_redirects=True) as client:
            for q in queries:
                if not q or len(q) < 3:
                    continue
                try:
                    url = f"https://en.wikipedia.org/w/api.php?action=query&generator=search&gsrsearch={urllib.parse.quote(q)}&gsrlimit=3&prop=pageimages&piprop=thumbnail&pithumbsize=800&format=json"
                    res = await client.get(url)
                    if res.status_code == 200:
                        data = res.json()
                        pages = data.get("query", {}).get("pages", {})
                        for _, p in pages.items():
                            thumb = p.get("thumbnail", {}).get("source")
                            if thumb and not thumb.lower().endswith(".svg"):
                                _REAL_IMAGE_URL_CACHE[cache_key] = thumb
                                return thumb
                            elif thumb:
                                _REAL_IMAGE_URL_CACHE[cache_key] = thumb
                                return thumb
                except Exception:
                    continue
    except Exception as e:
        logger.warning(f"Error querying real educational image for {keyword}: {e}")

    # Fallback to Pollinations AI real topic diagram
    import urllib.parse
    clean_title = re.sub(r'[^a-zA-Z0-9 ]', '', slide_title)[:50].strip()
    clean_top = re.sub(r'[^a-zA-Z0-9 ]', '', topic)[:40].strip()
    safe_prompt = urllib.parse.quote(f"clear educational illustration diagram of {clean_title} for {clean_top}, detailed science textbook quality, 8k")
    pollination_url = f"https://image.pollinations.ai/prompt/{safe_prompt}?width=800&height=500&nologo=true"
    _REAL_IMAGE_URL_CACHE[cache_key] = pollination_url
    return pollination_url


def _get_image_for_keyword(keyword: str) -> str:
    """Provides authentic educational image URL based on keyword theme without generic stock icons."""
    kw = (keyword or "").lower()
    # Neuroscience & Human Nervous System
    if any(k in kw for k in ["nervous", "brain", "neuron", "synapse", "spine", "nerve", "reflex", "cortex"]):
        return "https://images.unsplash.com/photo-1559757175-5700dde675bc?w=800&auto=format&fit=crop&q=80"
    # Cardiology & Circulatory System
    if any(k in kw for k in ["heart", "cardio", "blood", "circulat", "artery", "vessel"]):
        return "https://images.unsplash.com/photo-1530026405186-ed1f139313f8?w=800&auto=format&fit=crop&q=80"
    # Anatomy & Medical Biology
    if any(k in kw for k in ["anatomy", "skeleton", "body", "muscle", "organ", "digest", "respirat", "medical", "disease", "health"]):
        return "https://images.unsplash.com/photo-1576086213369-97a306d36557?w=800&auto=format&fit=crop&q=80"
    # Space & Astronomy
    if any(k in kw for k in ["space", "astronomy", "planet", "galaxy", "solar", "orbit", "universe", "star"]):
        return "https://images.unsplash.com/photo-1451187580459-43490279c0fa?w=800&auto=format&fit=crop&q=80"
    # Biology & Genetics
    if any(k in kw for k in ["bio", "dna", "cell", "organism", "genet", "microbe"]):
        return "https://images.unsplash.com/photo-1530026405186-ed1f139313f8?w=800&auto=format&fit=crop&q=80"
    # Botany & Plant Sciences
    if any(k in kw for k in ["plant", "photosynth", "flower", "leaf", "chloroplast", "botany", "crop", "forest"]):
        return "https://images.unsplash.com/photo-1518531933037-91b2f5f229cc?w=800&auto=format&fit=crop&q=80"
    # Chemistry & Lab Experiments
    if any(k in kw for k in ["chem", "molecule", "reaction", "lab", "experiment", "acid", "compound", "periodic", "element"]):
        return "https://images.unsplash.com/photo-1532094349884-543bc11b234d?w=800&auto=format&fit=crop&q=80"
    # Physics & Energy
    if any(k in kw for k in ["phys", "quantum", "electric", "magnet", "energy", "wave", "motion", "gravity", "optics"]):
        return "https://images.unsplash.com/photo-1507413245164-6160d8298b31?w=800&auto=format&fit=crop&q=80"
    # Mathematics & Geometry
    if any(k in kw for k in ["math", "geometry", "calculus", "algebra", "number", "vedic", "trig", "statistic"]):
        return "https://images.unsplash.com/photo-1635070041078-e363dbe005cb?w=800&auto=format&fit=crop&q=80"
    # History & Civilization
    if any(k in kw for k in ["history", "war", "revolut", "ancient", "monument", "civil", "india", "gandhi", "mughal", "empire"]):
        return "https://images.unsplash.com/photo-1461360370896-922624d12aa1?w=800&auto=format&fit=crop&q=80"
    # Technology & Computer Science
    if any(k in kw for k in ["ai", "robot", "comput", "tech", "program", "code", "cyber", "software", "network"]):
        return "https://images.unsplash.com/photo-1485827404703-89b55fcc595e?w=800&auto=format&fit=crop&q=80"
    # Earth Sciences, Climate & Geography
    if any(k in kw for k in ["earth", "climate", "environment", "geography", "eco", "river", "mountain", "soil", "ocean", "weather"]):
        return "https://images.unsplash.com/photo-1470071459604-3b5ec3a7fe05?w=800&auto=format&fit=crop&q=80"
    # Civics, Law & Politics
    if any(k in kw for k in ["civics", "polity", "constitution", "democracy", "parliament", "law", "government", "rights"]):
        return "https://images.unsplash.com/photo-1589829545856-d10d557cf95f?w=800&auto=format&fit=crop&q=80"
    # Literature & Languages
    if any(k in kw for k in ["liter", "english", "poem", "book", "lang", "grammar", "poetry", "novel"]):
        return "https://images.unsplash.com/photo-1497633762265-9d179a990aa6?w=800&auto=format&fit=crop&q=80"
    # Economics & Commerce
    if any(k in kw for k in ["econ", "market", "trade", "finance", "money", "commerce", "budget", "bank"]):
        return "https://images.unsplash.com/photo-1611974789855-9c2a0a7236a3?w=800&auto=format&fit=crop&q=80"
    # Sleek modern academic study & research background (NEVER an apple-on-books)
    return "https://images.unsplash.com/photo-1434030216411-0b793f4b4173?w=800&auto=format&fit=crop&q=80"


class PPTGeneratorService:

    async def generate_presentation(self, req: GeneratePPTRequest) -> PresentationData:
        """Generates a complete, high-quality presentation structure on any topic with teacher AI guidance."""
        presenter = (req.presenter_name or "").strip()
        if not presenter and req.user_email:
            presenter = req.user_email.split("@")[0].replace(".", " ").title()
        if not presenter:
            presenter = "Educator"

        prompt = f"""You are DEVGYA's Master Educational Presentation Architect.
Synthesize a comprehensive, beautifully structured slide deck for teachers and learners.

TOPIC: {req.topic}
TARGET AUDIENCE: {req.target_audience}
NUMBER OF SLIDES: {req.num_slides}
PRESENTATION TONE: {req.tone}
LANGUAGE: {req.language}
COLOR THEME: {req.theme}
PRESENTER / AUTHOR: {presenter}
{f"TEACHER SPECIFIC GUIDANCE & PEDAGOGY MANDATES: {req.teacher_guidance}" if req.teacher_guidance else ""}

CRITICAL ARCHITECTURAL MANDATES:
1. FIRST SLIDE (SLIDE 1) MUST BE A DEDICATED COVER / TITLE SLIDE:
   - "layout": "cover"
   - "category": "Presentation Cover"
   - "title": Main punchy presentation title
   - "subtitle": Clear, engaging subtitle summarizing audience goals
   - "bullets": [
       "Presented by: {presenter}",
       "Target Audience: {req.target_audience}",
       "Curriculum Focus: {req.topic}"
     ]
   - "speaker_notes": Welcoming opening script introducing the session, presenter, and overarching goals.

2. FINAL SLIDE (SLIDE {req.num_slides}) MUST BE A DEDICATED THANK YOU & DISCUSSION SLIDE:
   - "layout": "thank_you"
   - "category": "Conclusion & Discussion"
   - "title": "Thank You!"
   - "subtitle": "Questions & Classroom Discussion"
   - "bullets": [
       "**Core Key Takeaway**: [1 crisp, memorable summary sentence]",
       "**Classroom Discussion Question**: [1 thought-provoking discussion prompt for students]",
       "**Next Steps & Review**: Concept consolidation, chapter exercises, and open Q&A"
     ]
   - "speaker_notes": Warm closing remarks thanking students/audience and opening the floor for discussion.

3. DYNAMIC & FRESH LOOKS (NO REPETITIVE MONOTONOUS SLIDES):
   Vary slide layouts across middle slides (Slides 2 to {req.num_slides - 1}) based on the specific content:
   - 'two_column': Comparison, theoretical vs practical, advantages vs challenges.
   - 'stat_highlight': 2-3 prominent quantitative metrics or pivotal numbers with values and descriptions.
   - 'process_timeline': 3-4 sequential stages, milestones, or procedural steps with titles and descriptions.
   - 'quote_insight': Powerful conceptual quote, foundational axiom, or thought leader insight.
   - 'split_image_text': High-impact concept explanation paired with visual focal illustration.
   - 'title_bullets': Structured points with bold lead-in keywords (**Concept**: Explanation).
   Every slide must feel intentionally crafted, professional, and visually distinct. Never repeat identical layout formats consecutively.

RETURN VALID JSON ONLY matching this exact schema:
{{
  "title": "Main Presentation Title",
  "subtitle": "Clear, engaging subtitle summarizing audience goal",
  "slides": [
    {{
      "slide_number": 1,
      "layout": "cover",
      "category": "Presentation Cover",
      "title": "{req.topic}",
      "subtitle": "A Comprehensive Guide for {req.target_audience}",
      "bullets": [
        "Presented by: {presenter}",
        "Target Audience: {req.target_audience}",
        "Subject: {req.topic}"
      ],
      "left_column": null,
      "right_column": null,
      "metrics": null,
      "timeline_steps": null,
      "quote": null,
      "image_keyword": "{req.topic}",
      "image_caption": "Presentation Cover",
      "speaker_notes": "Welcome everyone to today's session on {req.topic}..."
    }},
    {{
      "slide_number": 2,
      "layout": "two_column",
      "category": "Core Mechanism Comparison",
      "title": "...",
      "subtitle": "...",
      "bullets": [],
      "left_column": {{
        "title": "Classical View",
        "bullets": ["Point A", "Point B"]
      }},
      "right_column": {{
        "title": "Modern View",
        "bullets": ["Point X", "Point Y"]
      }},
      "metrics": null,
      "timeline_steps": null,
      "quote": null,
      "image_keyword": "laboratory science",
      "image_caption": "Comparative analysis",
      "speaker_notes": "Highlight how the transition occurred..."
    }}
  ]
}}

Generate ALL {req.num_slides} slides completely!"""

        try:
            raw = await ai_provider.chat_completion(
                messages=[
                    {"role": "system", "content": "You are DEVGYA's premier AI Slide Deck Architect. Output strictly valid JSON without markdown formatting."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.35,
                max_tokens=4000,
                response_format_json=True
            )

            parsed = robust_json_parser(raw)
            slides_raw = parsed.get("slides") or []

            # If empty or malformed fallback
            if not slides_raw or not isinstance(slides_raw, list):
                raise ValueError("No valid slides parsed from AI response.")

            # Concurrently resolve real educational topic images for every slide
            async def _fill_slide_image(s_dict):
                if not isinstance(s_dict, dict):
                    return None
                kw = str(s_dict.get("image_keyword") or req.topic)
                stitle = str(s_dict.get("title") or req.topic)
                img = s_dict.get("image_url")
                if not img or "unsplash.com" in str(img):
                    img = await _resolve_real_topic_image(req.topic, stitle, kw)
                return img

            resolved_imgs = await asyncio.gather(*[_fill_slide_image(s) for s in slides_raw], return_exceptions=True)

            slides_list: List[SlideItem] = []
            total_s = len(slides_raw)
            for idx, s in enumerate(slides_raw):
                if not isinstance(s, dict):
                    continue
                num = idx + 1
                kw = str(s.get("image_keyword") or req.topic)
                r_img = resolved_imgs[idx] if idx < len(resolved_imgs) and isinstance(resolved_imgs[idx], str) and resolved_imgs[idx] else None
                img_url = r_img or s.get("image_url") or _get_image_for_keyword(kw)

                bullets = s.get("bullets") if isinstance(s.get("bullets"), list) else []
                if not bullets and s.get("content"):
                    bullets = [str(s.get("content"))]

                layout_str = str(s.get("layout") or "title_bullets")
                cat_str = str(s.get("category") or f"Module {num}")
                title_str = str(s.get("title") or f"Key Concept {num}")
                sub_str = str(s.get("subtitle") or "") if s.get("subtitle") else None

                # Enforce Slide 1 Cover & Final Slide Thank You
                if num == 1:
                    layout_str = "cover"
                    cat_str = "Presentation Cover"
                    if not any("Presented by" in b for b in bullets):
                        bullets = [f"Presented by: {presenter}"] + [b for b in bullets if "Presented by" not in b]
                elif num == total_s:
                    layout_str = "thank_you"
                    cat_str = "Conclusion & Discussion"
                    if "thank" not in title_str.lower():
                        title_str = "Thank You!"
                    if not sub_str:
                        sub_str = "Questions & Classroom Discussion"

                item = SlideItem(
                    slide_number=num,
                    layout=layout_str,
                    category=cat_str,
                    title=title_str,
                    subtitle=sub_str,
                    bullets=bullets,
                    left_column=s.get("left_column") if isinstance(s.get("left_column"), dict) else None,
                    right_column=s.get("right_column") if isinstance(s.get("right_column"), dict) else None,
                    metrics=s.get("metrics") if isinstance(s.get("metrics"), list) else None,
                    timeline_steps=s.get("timeline_steps") if isinstance(s.get("timeline_steps"), list) else None,
                    quote=s.get("quote") if isinstance(s.get("quote"), dict) else None,
                    image_keyword=kw,
                    image_url=img_url,
                    image_caption=s.get("image_caption") or f"Visual guide for {req.topic}",
                    speaker_notes=str(s.get("speaker_notes") or f"Guide students through key ideas of this slide.")
                )
                slides_list.append(item)

            if not slides_list:
                raise ValueError("Could not assemble valid slides from LLM output.")

            deck_res = PresentationData(
                id=f"ppt-{uuid.uuid4().hex[:12]}",
                title=str(parsed.get("title") or req.topic),
                subtitle=str(parsed.get("subtitle") or f"A comprehensive study presentation for {req.target_audience}"),
                topic=req.topic,
                target_audience=req.target_audience,
                num_slides=len(slides_list),
                theme=req.theme,
                language=req.language,
                teacher_guidance=req.teacher_guidance,
                presenter_name=presenter,
                slides=slides_list
            )

            # Auto-save deck into user's personal Supabase cloud history
            clean_user = (req.user_id or req.user_email or "").strip()
            if clean_user:
                try:
                    ppt_history_service.save_deck(clean_user, deck_res.dict())
                except Exception as save_err:
                    logger.warning(f"Failed to auto-save deck to Supabase: {save_err}")

            return deck_res

        except Exception as e:
            logger.warning(f"AI presentation generation error: {e}. Generating fallback structured presentation.")
            return self._generate_fallback_presentation(req)

    async def resolve_real_image(self, topic: str, slide_title: str, keyword: str) -> str:
        """Finds a real educational image for a topic or slide."""
        return await _resolve_real_topic_image(topic, slide_title, keyword)

    def _generate_fallback_presentation(self, req: GeneratePPTRequest) -> PresentationData:
        """Fallback presentation structure with dedicated Cover and Thank You slides."""
        presenter = (req.presenter_name or "").strip()
        if not presenter and req.user_email:
            presenter = req.user_email.split("@")[0].replace(".", " ").title()
        if not presenter:
            presenter = "Educator"
        topic_title = req.topic.strip().title()

        slides: List[SlideItem] = [
            SlideItem(
                slide_number=1,
                layout="cover",
                category="Presentation Cover",
                title=topic_title,
                subtitle=f"A Comprehensive Pedagogical Guide for {req.target_audience}",
                bullets=[
                    f"Presented by: {presenter}",
                    f"Target Audience: {req.target_audience}",
                    f"Curriculum Focus: {topic_title}"
                ],
                image_keyword=req.topic,
                image_url=_get_image_for_keyword(req.topic),
                image_caption=f"Overview of {topic_title}",
                speaker_notes=f"Welcome students to the session on {topic_title}. Outline the central learning objectives."
            ),
            SlideItem(
                slide_number=2,
                layout="two_column",
                category="Key Fundamentals",
                title="Foundational Principles & Concepts",
                subtitle="Understanding core mechanisms",
                left_column={
                    "title": "Theoretical Framework",
                    "bullets": [
                        f"Fundamental definitions governing {topic_title}.",
                        "Standard scientific and academic axioms.",
                        "Direct connection to NCERT/CBSE benchmarks."
                    ]
                },
                right_column={
                    "title": "Practical Application",
                    "bullets": [
                        "Real-world observation and laboratory relevance.",
                        "Everyday case scenarios and contextual examples.",
                        "Common problem-solving methodologies."
                    ]
                },
                image_keyword="science research study",
                image_url=_get_image_for_keyword("science research"),
                image_caption="Theoretical vs Practical Dimensions",
                speaker_notes="Walk students through the key distinction between foundational theory and observable applications."
            ),
            SlideItem(
                slide_number=3,
                layout="stat_highlight",
                category="Analysis & Metrics",
                title="Key Metrics & Significance",
                subtitle="Quantitative and analytical dimensions",
                metrics=[
                    {"label": "Core Impact", "value": "100%", "description": f"Essential mastery for {req.target_audience}."},
                    {"label": "Retention Rate", "value": "85%+", "description": "Achieved via structured conceptual visualization."},
                    {"label": "Exam Weightage", "value": "High", "description": "Frequently featured in standard examination blueprints."}
                ],
                image_keyword="analytics data chart",
                image_url=_get_image_for_keyword("data chart"),
                image_caption="Analytical Framework",
                speaker_notes="Emphasize why this topic carries significant importance in academic evaluation."
            )
        ]

        # Add middle slides if needed
        while len(slides) < req.num_slides - 1:
            num = len(slides) + 1
            slides.append(
                SlideItem(
                    slide_number=num,
                    layout="title_bullets",
                    category=f"Module {num}",
                    title=f"{topic_title} — Analytical Deep Dive Part {num}",
                    subtitle="Detailed exploration of advanced implications",
                    bullets=[
                        f"**In-Depth Aspect {num}.1**: Critical exploration of underlying concepts.",
                        f"**Practical Example**: Step-by-step examination of common problem patterns.",
                        f"**Misconception Alert**: Addressing common learner pitfalls."
                    ],
                    image_keyword=req.topic,
                    image_url=_get_image_for_keyword(req.topic),
                    image_caption=f"Deep-dive analysis of {topic_title}",
                    speaker_notes=f"Deep dive into Module {num}. Encourage students to take concise notes."
                )
            )

        # Final Slide: Thank You
        final_num = len(slides) + 1
        slides.append(
            SlideItem(
                slide_number=final_num,
                layout="thank_you",
                category="Conclusion & Discussion",
                title="Thank You!",
                subtitle="Questions & Classroom Discussion",
                bullets=[
                    f"**Core Takeaway**: Essential conceptual mastery of {topic_title}.",
                    "**Classroom Discussion Question**: How do these concepts impact modern practice and future developments?",
                    "**Next Steps**: Review notes, chapter questions, and collaborative discussion."
                ],
                image_keyword="classroom celebration education",
                image_url=_get_image_for_keyword("classroom discussion"),
                image_caption="Questions and Discussion",
                speaker_notes="Thank everyone for their attention and open the floor for questions."
            )
        )

        fallback_deck = PresentationData(
            id=f"ppt-{uuid.uuid4().hex[:12]}",
            title=topic_title,
            subtitle=f"Comprehensive study presentation for {req.target_audience}",
            topic=req.topic,
            target_audience=req.target_audience,
            num_slides=len(slides),
            theme=req.theme,
            language=req.language,
            teacher_guidance=req.teacher_guidance,
            presenter_name=presenter,
            slides=slides
        )

        clean_user = (req.user_id or req.user_email or "").strip()
        if clean_user:
            try:
                ppt_history_service.save_deck(clean_user, fallback_deck.dict())
            except Exception as save_err:
                logger.warning(f"Failed to auto-save fallback deck to Supabase: {save_err}")

        return fallback_deck

        clean_user = (req.user_id or req.user_email or "").strip()
        if clean_user:
            try:
                ppt_history_service.save_deck(clean_user, fallback_deck.dict())
            except Exception as save_err:
                logger.warning(f"Failed to auto-save fallback deck to Supabase: {save_err}")

        return fallback_deck

    async def refine_slide(self, req: RefineSlideRequest) -> SlideItem:
        """Refines or rephrases an individual slide based on teacher's specific instruction."""
        prompt = f"""You are DEVGYA's AI Slide Polisher.
Refine this educational slide based on the teacher's instruction.

TOPIC: {req.topic}
TARGET AUDIENCE: {req.target_audience}
TEACHER INSTRUCTION: {req.instruction}

CURRENT SLIDE DATA:
{json.dumps(req.slide.dict(), indent=2)}

RETURN UPDATED SLIDE JSON ONLY with the same keys (slide_number, layout, category, title, subtitle, bullets, left_column, right_column, metrics, timeline_steps, quote, image_keyword, image_url, image_caption, speaker_notes)."""

        try:
            raw = await ai_provider.chat_completion(
                messages=[
                    {"role": "system", "content": "You are DEVGYA's Slide Refiner. Return valid JSON only."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.3,
                max_tokens=2000,
                response_format_json=True
            )
            parsed = robust_json_parser(raw)
            kw = str(parsed.get("image_keyword") or req.slide.image_keyword)
            return SlideItem(
                slide_number=req.slide.slide_number,
                layout=str(parsed.get("layout") or req.slide.layout),
                category=str(parsed.get("category") or req.slide.category),
                title=str(parsed.get("title") or req.slide.title),
                subtitle=str(parsed.get("subtitle") or req.slide.subtitle or ""),
                bullets=parsed.get("bullets") or req.slide.bullets,
                left_column=parsed.get("left_column") or req.slide.left_column,
                right_column=parsed.get("right_column") or req.slide.right_column,
                metrics=parsed.get("metrics") or req.slide.metrics,
                timeline_steps=parsed.get("timeline_steps") or req.slide.timeline_steps,
                quote=parsed.get("quote") or req.slide.quote,
                image_keyword=kw,
                image_url=parsed.get("image_url") or req.slide.image_url or _get_image_for_keyword(kw),
                image_caption=parsed.get("image_caption") or req.slide.image_caption,
                speaker_notes=str(parsed.get("speaker_notes") or req.slide.speaker_notes)
            )
        except Exception as e:
            status_code, detail = format_ai_exception_detail(e, "Slide Polish")
            raise HTTPException(status_code=status_code, detail=detail)

    def generate_pptx(self, pres_data: PresentationData) -> bytes:
        """Builds a real 16:9 Microsoft PowerPoint (.pptx) file."""
        prs = Presentation()
        # Set 16:9 Widescreen aspect ratio
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6] # Blank slide
        theme = THEME_PRESETS.get(pres_data.theme, THEME_PRESETS["modern_navy"])
        c_primary = RGBColor(*theme["rgb_primary"])
        c_accent = RGBColor(*theme["rgb_accent"])
        c_card = RGBColor(*theme["rgb_card"])
        c_bg = RGBColor(*theme["rgb_bg"])
        presenter_display = pres_data.presenter_name or "Educator"

        for s in pres_data.slides:
            slide = prs.slides.add_slide(blank_layout)

            # Slide background card shape
            bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = c_bg
            bg_shape.line.fill.background()

            # 1. SPECIAL COVER SLIDE (SLIDE 1)
            if s.layout == "cover" or s.slide_number == 1:
                # Left decorative vertical accent banner
                left_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
                left_band.fill.solid()
                left_band.fill.fore_color.rgb = c_accent
                left_band.line.fill.background()

                # Brand Pill
                pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.0), Inches(5.0), Inches(0.45))
                pill.fill.solid()
                pill.fill.fore_color.rgb = c_card
                pill.line.color.rgb = c_accent
                tf_p = pill.text_frame
                tf_p.vertical_anchor = MSO_ANCHOR.MIDDLE
                p_pill = tf_p.paragraphs[0]
                p_pill.text = "DEVGYA AI  •  EDUCATIONAL PRESENTATION"
                p_pill.font.size = Pt(10)
                p_pill.font.bold = True
                p_pill.font.color.rgb = c_accent

                # Main Topic Title
                title_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(11.0), Inches(2.3))
                tf_title = title_box.text_frame
                tf_title.word_wrap = True
                p_t = tf_title.paragraphs[0]
                p_t.text = s.title
                p_t.font.size = Pt(36)
                p_t.font.bold = True
                p_t.font.color.rgb = c_primary

                if s.subtitle:
                    p_sub = tf_title.add_paragraph()
                    p_sub.text = s.subtitle
                    p_sub.font.size = Pt(18)
                    p_sub.font.color.rgb = RGBColor(71, 85, 105)
                    p_sub.space_before = Pt(10)

                # Presenter & Attribution Card
                pres_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(4.5), Inches(10.8), Inches(1.8))
                pres_card.fill.solid()
                pres_card.fill.fore_color.rgb = c_card
                pres_card.line.color.rgb = RGBColor(226, 232, 240)
                tf_pc = pres_card.text_frame
                tf_pc.margin_left = Inches(0.4)
                tf_pc.margin_top = Inches(0.25)

                p_by = tf_pc.paragraphs[0]
                p_by.text = f"Presented by: {presenter_display}"
                p_by.font.size = Pt(16)
                p_by.font.bold = True
                p_by.font.color.rgb = c_primary

                p_aud = tf_pc.add_paragraph()
                p_aud.text = f"Target Audience: {pres_data.target_audience}   •   Subject Focus: {pres_data.topic}"
                p_aud.font.size = Pt(12)
                p_aud.font.color.rgb = RGBColor(100, 116, 139)
                p_aud.space_before = Pt(6)

                p_org = tf_pc.add_paragraph()
                p_org.text = "DEVGYA Global Edutech Private Limited  •  CBSE & NCERT Aligned"
                p_org.font.size = Pt(11)
                p_org.font.bold = True
                p_org.font.color.rgb = c_accent
                p_org.space_before = Pt(6)

                if s.speaker_notes:
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = f"TEACHER SPEAKER NOTES:\n{s.speaker_notes}"
                continue

            # 2. SPECIAL THANK YOU SLIDE (FINAL SLIDE)
            if s.layout == "thank_you" or s.slide_number == len(pres_data.slides):
                ty_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(0.8), Inches(10.933), Inches(5.9))
                ty_card.fill.solid()
                ty_card.fill.fore_color.rgb = c_card
                ty_card.line.color.rgb = c_accent

                tf_ty = ty_card.text_frame
                tf_ty.margin_left = Inches(0.6)
                tf_ty.margin_right = Inches(0.6)
                tf_ty.margin_top = Inches(0.4)

                p_th = tf_ty.paragraphs[0]
                p_th.alignment = PP_ALIGN.CENTER
                p_th.text = "Thank You!"
                p_th.font.size = Pt(40)
                p_th.font.bold = True
                p_th.font.color.rgb = c_primary

                p_tsub = tf_ty.add_paragraph()
                p_tsub.alignment = PP_ALIGN.CENTER
                p_tsub.text = s.subtitle or "Questions & Classroom Discussion"
                p_tsub.font.size = Pt(18)
                p_tsub.font.bold = True
                p_tsub.font.color.rgb = c_accent
                p_tsub.space_before = Pt(6)

                for b in s.bullets:
                    clean_b = b.replace("**", "").replace("*", "")
                    p_b = tf_ty.add_paragraph()
                    p_b.text = f"• {clean_b}"
                    p_b.font.size = Pt(14)
                    p_b.font.color.rgb = RGBColor(51, 65, 85)
                    p_b.space_before = Pt(12)

                p_sig = tf_ty.add_paragraph()
                p_sig.alignment = PP_ALIGN.CENTER
                p_sig.text = f"Presented by: {presenter_display}   •   DEVGYA Global Edutech"
                p_sig.font.size = Pt(12)
                p_sig.font.bold = True
                p_sig.font.color.rgb = RGBColor(148, 163, 184)
                p_sig.space_before = Pt(24)

                if s.speaker_notes:
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = f"TEACHER SPEAKER NOTES:\n{s.speaker_notes}"
                continue

            # 3. MIDDLE SLIDES WITH TOP ACCENT HEADER BAR
            header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.12))
            header_bar.fill.solid()
            header_bar.fill.fore_color.rgb = c_accent
            header_bar.line.fill.background()

            # Category Pill & Slide Number
            cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(9.0), Inches(0.4))
            tf_cat = cat_box.text_frame
            tf_cat.word_wrap = True
            p_cat = tf_cat.paragraphs[0]
            p_cat.text = f"{s.category.upper()}  •  DEVGYA AI STUDY SUITE"
            p_cat.font.size = Pt(10)
            p_cat.font.bold = True
            p_cat.font.color.rgb = c_accent

            num_box = slide.shapes.add_textbox(Inches(10.5), Inches(0.7), Inches(2.0), Inches(0.4))
            p_num = num_box.text_frame.paragraphs[0]
            p_num.alignment = PP_ALIGN.RIGHT
            p_num.text = f"Slide {s.slide_number} of {len(pres_data.slides)}"
            p_num.font.size = Pt(10)
            p_num.font.bold = True
            p_num.font.color.rgb = RGBColor(148, 163, 184)

            # Main Slide Title & Subtitle
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.733), Inches(1.1))
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            p_t = tf_title.paragraphs[0]
            p_t.text = s.title
            p_t.font.size = Pt(24)
            p_t.font.bold = True
            p_t.font.color.rgb = c_primary

            if s.subtitle:
                p_sub = tf_title.add_paragraph()
                p_sub.text = s.subtitle
                p_sub.font.size = Pt(13)
                p_sub.font.color.rgb = RGBColor(71, 85, 105)

            # Layout specific body rendering
            content_top = Inches(2.3)
            content_height = Inches(4.3)

            if s.layout == "two_column" and (s.left_column or s.right_column):
                col_w = Inches(5.6)
                card_l = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), content_top, col_w, content_height)
                card_l.fill.solid()
                card_l.fill.fore_color.rgb = c_card
                card_l.line.color.rgb = RGBColor(226, 232, 240)

                tf_l = card_l.text_frame
                tf_l.word_wrap = True
                tf_l.margin_left = Inches(0.3)
                tf_l.margin_right = Inches(0.3)
                tf_l.margin_top = Inches(0.3)

                col_l_title = (s.left_column or {}).get("title", "Part 1")
                p_lt = tf_l.paragraphs[0]
                p_lt.text = col_l_title
                p_lt.font.size = Pt(16)
                p_lt.font.bold = True
                p_lt.font.color.rgb = c_primary

                for b in (s.left_column or {}).get("bullets", []):
                    p_b = tf_l.add_paragraph()
                    p_b.text = f"• {b}"
                    p_b.font.size = Pt(12)
                    p_b.font.color.rgb = RGBColor(51, 65, 85)

                card_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), content_top, col_w, content_height)
                card_r.fill.solid()
                card_r.fill.fore_color.rgb = c_card
                card_r.line.color.rgb = RGBColor(226, 232, 240)

                tf_r = card_r.text_frame
                tf_r.word_wrap = True
                tf_r.margin_left = Inches(0.3)
                tf_r.margin_right = Inches(0.3)
                tf_r.margin_top = Inches(0.3)

                col_r_title = (s.right_column or {}).get("title", "Part 2")
                p_rt = tf_r.paragraphs[0]
                p_rt.text = col_r_title
                p_rt.font.size = Pt(16)
                p_rt.font.bold = True
                p_rt.font.color.rgb = c_primary

                for b in (s.right_column or {}).get("bullets", []):
                    p_b = tf_r.add_paragraph()
                    p_b.text = f"• {b}"
                    p_b.font.size = Pt(12)
                    p_b.font.color.rgb = RGBColor(51, 65, 85)

            elif s.layout == "stat_highlight" and s.metrics:
                num_stats = min(len(s.metrics), 4)
                box_w = Inches(11.733 / num_stats - 0.2)
                for i, m in enumerate(s.metrics[:num_stats]):
                    left_pos = Inches(0.8 + i * (11.733 / num_stats))
                    stat_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, content_top, box_w, content_height)
                    stat_card.fill.solid()
                    stat_card.fill.fore_color.rgb = c_card
                    stat_card.line.color.rgb = c_accent

                    tf_s = stat_card.text_frame
                    tf_s.word_wrap = True
                    tf_s.vertical_anchor = MSO_ANCHOR.MIDDLE

                    p_val = tf_s.paragraphs[0]
                    p_val.text = str(m.get("value", ""))
                    p_val.font.size = Pt(32)
                    p_val.font.bold = True
                    p_val.font.color.rgb = c_accent
                    p_val.alignment = PP_ALIGN.CENTER

                    p_lbl = tf_s.add_paragraph()
                    p_lbl.text = str(m.get("label", ""))
                    p_lbl.font.size = Pt(13)
                    p_lbl.font.bold = True
                    p_lbl.font.color.rgb = c_primary
                    p_lbl.alignment = PP_ALIGN.CENTER

            elif s.layout == "quote_insight" and s.quote:
                q_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), content_top, Inches(10.333), content_height)
                q_card.fill.solid()
                q_card.fill.fore_color.rgb = c_card
                q_card.line.color.rgb = c_accent

                tf_q = q_card.text_frame
                tf_q.word_wrap = True
                tf_q.margin_left = Inches(0.6)
                tf_q.margin_right = Inches(0.6)
                tf_q.vertical_anchor = MSO_ANCHOR.MIDDLE

                p_q = tf_q.paragraphs[0]
                p_q.text = f"“{s.quote.get('text', '')}”"
                p_q.font.size = Pt(20)
                p_q.font.italic = True
                p_q.font.color.rgb = c_primary
                p_q.alignment = PP_ALIGN.CENTER

                p_qa = tf_q.add_paragraph()
                p_qa.text = f"— {s.quote.get('author', 'Core Pedagogical Principle')}"
                p_qa.font.size = Pt(13)
                p_qa.font.bold = True
                p_qa.font.color.rgb = c_accent
                p_qa.alignment = PP_ALIGN.CENTER

            elif s.layout == "process_timeline" and s.timeline_steps:
                num_steps = min(len(s.timeline_steps), 4)
                step_w = Inches(11.733 / num_steps - 0.2)
                for i, st in enumerate(s.timeline_steps[:num_steps]):
                    left_pos = Inches(0.8 + i * (11.733 / num_steps))
                    step_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, content_top, step_w, content_height)
                    step_card.fill.solid()
                    step_card.fill.fore_color.rgb = c_card
                    step_card.line.color.rgb = c_accent

                    tf_step = step_card.text_frame
                    tf_step.word_wrap = True
                    tf_step.margin_left = Inches(0.25)
                    tf_step.margin_right = Inches(0.25)
                    tf_step.margin_top = Inches(0.4)

                    p_step_num = tf_step.paragraphs[0]
                    p_step_num.text = f"STEP {i+1}"
                    p_step_num.font.size = Pt(11)
                    p_step_num.font.bold = True
                    p_step_num.font.color.rgb = c_accent

                    p_step_title = tf_step.add_paragraph()
                    p_step_title.text = str(st.get("title") or st.get("step") or f"Phase {i+1}")
                    p_step_title.font.size = Pt(14)
                    p_step_title.font.bold = True
                    p_step_title.font.color.rgb = c_primary
                    p_step_title.space_before = Pt(6)

                    p_step_desc = tf_step.add_paragraph()
                    p_step_desc.text = str(st.get("description") or "")
                    p_step_desc.font.size = Pt(11)
                    p_step_desc.font.color.rgb = RGBColor(71, 85, 105)
                    p_step_desc.space_before = Pt(8)

            else:
                card_bullets = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), content_top, Inches(7.5), content_height)
                card_bullets.fill.solid()
                card_bullets.fill.fore_color.rgb = c_card
                card_bullets.line.color.rgb = RGBColor(226, 232, 240)

                tf_b = card_bullets.text_frame
                tf_b.word_wrap = True
                tf_b.margin_left = Inches(0.4)
                tf_b.margin_right = Inches(0.4)
                tf_b.margin_top = Inches(0.4)

                first_p = True
                for b in s.bullets:
                    clean_b = b.replace("**", "").replace("*", "")
                    if first_p:
                        p_line = tf_b.paragraphs[0]
                        first_p = False
                    else:
                        p_line = tf_b.add_paragraph()
                    p_line.text = f"• {clean_b}"
                    p_line.font.size = Pt(13)
                    p_line.font.color.rgb = RGBColor(30, 41, 59)
                    p_line.space_after = Pt(10)

                img_data = _download_image_bytes(s.image_url) if s.image_url else None
                pic_inserted = False
                if img_data:
                    try:
                        slide.shapes.add_picture(
                            io.BytesIO(img_data),
                            Inches(8.6), content_top,
                            Inches(3.9), content_height - Inches(0.75)
                        )
                        cap_shape = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(8.6), content_top + content_height - Inches(0.65),
                            Inches(3.9), Inches(0.65)
                        )
                        cap_shape.fill.solid()
                        cap_shape.fill.fore_color.rgb = c_card
                        cap_shape.line.color.rgb = RGBColor(226, 232, 240)
                        tf_c = cap_shape.text_frame
                        tf_c.word_wrap = True
                        tf_c.margin_top = Inches(0.08)
                        tf_c.margin_left = Inches(0.15)
                        p_cap = tf_c.paragraphs[0]
                        p_cap.text = s.image_caption or s.image_keyword or "Visual Reference"
                        p_cap.font.size = Pt(10)
                        p_cap.font.color.rgb = RGBColor(51, 65, 85)
                        pic_inserted = True
                    except Exception as pic_err:
                        logger.warning(f"PPTX picture insertion failed: {pic_err}")
                        pic_inserted = False

                if not pic_inserted:
                    card_img = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.7), content_top, Inches(3.8), content_height)
                    card_img.fill.solid()
                    card_img.fill.fore_color.rgb = c_primary
                    card_img.line.fill.background()

                    tf_img = card_img.text_frame
                    tf_img.word_wrap = True
                    tf_img.margin_left = Inches(0.3)
                    tf_img.margin_right = Inches(0.3)
                    tf_img.vertical_anchor = MSO_ANCHOR.MIDDLE

                    p_ik = tf_img.paragraphs[0]
                    p_ik.text = f"VISUAL FOCUS"
                    p_ik.font.size = Pt(11)
                    p_ik.font.bold = True
                    p_ik.font.color.rgb = c_accent

                    p_cap = tf_img.add_paragraph()
                    p_cap.text = s.image_caption or s.image_keyword
                    p_cap.font.size = Pt(14)
                    p_cap.font.bold = True
                    p_cap.font.color.rgb = RGBColor(255, 255, 255)

                    p_link = tf_img.add_paragraph()
                    p_link.text = f"Topic: {pres_data.topic}\nAudience: {pres_data.target_audience}"
                    p_link.font.size = Pt(10)
                    p_link.font.color.rgb = RGBColor(203, 213, 225)

            # Speaker Notes in PPTX
            if s.speaker_notes:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = f"TEACHER SPEAKER NOTES:\n{s.speaker_notes}"

        out_buffer = io.BytesIO()
        prs.save(out_buffer)
        out_bytes = out_buffer.getvalue()
        out_buffer.close()
        return out_bytes

    def generate_pdf(self, pres_data: PresentationData) -> bytes:
        """Builds a high-impact landscape slide deck PDF using ReportLab."""
        buffer = io.BytesIO()

        # Landscape A4 size: 841.89 x 595.27 points
        doc = SimpleDocTemplate(
            buffer,
            pagesize=landscape(A4),
            leftMargin=36,
            rightMargin=36,
            topMargin=32,
            bottomMargin=32
        )

        theme = THEME_PRESETS.get(pres_data.theme, THEME_PRESETS["modern_navy"])
        styles = getSampleStyleSheet()

        title_style = ParagraphStyle(
            'SlideTitle',
            parent=styles['Heading1'],
            fontName=UNICODE_BOLD_FONT_NAME,
            fontSize=18,
            leading=22,
            textColor=colors.HexColor(theme["primary"])
        )

        subtitle_style = ParagraphStyle(
            'SlideSubtitle',
            parent=styles['Normal'],
            fontName=UNICODE_FONT_NAME,
            fontSize=11,
            leading=15,
            textColor=colors.HexColor(theme["text_secondary"])
        )

        cat_badge_style = ParagraphStyle(
            'CatBadge',
            parent=styles['Normal'],
            fontName=UNICODE_BOLD_FONT_NAME,
            fontSize=9,
            leading=12,
            textColor=colors.HexColor(theme["accent"])
        )

        bullet_style = ParagraphStyle(
            'BulletText',
            parent=styles['Normal'],
            fontName=UNICODE_FONT_NAME,
            fontSize=10.5,
            leading=16,
            textColor=colors.HexColor(theme["text_primary"])
        )

        notes_style = ParagraphStyle(
            'NotesText',
            parent=styles['Normal'],
            fontName=UNICODE_FONT_NAME,
            fontSize=8.5,
            leading=12,
            textColor=colors.HexColor("#475569")
        )

        story = []

        total_slides = len(pres_data.slides)
        presenter_display = pres_data.presenter_name or "Educator"

        for idx, s in enumerate(pres_data.slides):
            slide_elements = []

            # 1. SPECIAL COVER SLIDE (SLIDE 1)
            if s.layout == "cover" or s.slide_number == 1:
                # Top Pill & Brand
                top_cover_table = Table(
                    [[Paragraph("<b>DEVGYA AI  •  EDUCATIONAL PRESENTATION</b>", cat_badge_style), Paragraph("<b>CBSE & NCERT Aligned</b>", cat_badge_style)]],
                    colWidths=[550, 220]
                )
                top_cover_table.setStyle(TableStyle([
                    ('ALIGN', (0,0), (0,0), 'LEFT'),
                    ('ALIGN', (1,0), (1,0), 'RIGHT'),
                    ('PADDING', (0,0), (-1,-1), 0),
                    ('BOTTOMPADDING', (0,0), (-1,-1), 4),
                ]))
                slide_elements.append(top_cover_table)
                slide_elements.append(HRFlowable(width="100%", thickness=3, color=colors.HexColor(theme["accent"]), spaceBefore=2, spaceAfter=20))

                # Large Cover Title & Subtitle
                cover_title_style = ParagraphStyle(
                    'CoverTitle',
                    parent=title_style,
                    fontSize=28,
                    leading=34,
                    textColor=colors.HexColor(theme["primary"])
                )
                cover_sub_style = ParagraphStyle(
                    'CoverSub',
                    parent=subtitle_style,
                    fontSize=13,
                    leading=18,
                    textColor=colors.HexColor(theme["text_secondary"])
                )
                slide_elements.append(Paragraph(clean_md_to_reportlab(strip_emojis_for_pdf(s.title)), cover_title_style))
                if s.subtitle:
                    slide_elements.append(Spacer(1, 8))
                    slide_elements.append(Paragraph(clean_md_to_reportlab(strip_emojis_for_pdf(s.subtitle)), cover_sub_style))

                slide_elements.append(Spacer(1, 24))

                # Presenter Attribution Card
                pres_card_content = [
                    Paragraph(f"<font size=14 color='{theme['primary']}'><b>Presented by: {presenter_display}</b></font>", bullet_style),
                    Spacer(1, 4),
                    Paragraph(f"<b>Target Audience:</b> {pres_data.target_audience}   •   <b>Subject Focus:</b> {strip_emojis_for_pdf(pres_data.topic)}", subtitle_style),
                    Spacer(1, 4),
                    Paragraph(f"<b>DEVGYA Global Edutech Private Limited</b>", cat_badge_style)
                ]
                pres_table = Table([[pres_card_content]], colWidths=[760])
                pres_table.setStyle(TableStyle([
                    ('BACKGROUND', (0,0), (-1,-1), colors.HexColor(theme["card_bg"])),
                    ('BOX', (0,0), (-1,-1), 1.5, colors.HexColor(theme["accent"])),
                    ('PADDING', (0,0), (-1,-1), 14),
                ]))
                slide_elements.append(pres_table)

                if s.speaker_notes:
                    slide_elements.append(Spacer(1, 14))
                    notes_clean = clean_md_to_reportlab(strip_emojis_for_pdf(s.speaker_notes))
                    notes_table = Table([[Paragraph(f"<b>PRESENTER OPENING NOTES:</b> {notes_clean}", notes_style)]], colWidths=[760])
                    notes_table.setStyle(TableStyle([
                        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor("#F1F5F9")),
                        ('BOX', (0,0), (-1,-1), 0.5, colors.HexColor("#94A3B8")),
                        ('PADDING', (0,0), (-1,-1), 6),
                    ]))
                    slide_elements.append(notes_table)

                story.append(KeepTogether(slide_elements))
                if idx < total_slides - 1:
                    from reportlab.platypus import PageBreak
                    story.append(PageBreak())
                continue

            # 2. SPECIAL THANK YOU SLIDE (FINAL SLIDE)
            if s.layout == "thank_you" or s.slide_number == total_slides:
                ty_h_style = ParagraphStyle(
                    'TYTitle',
                    parent=title_style,
                    fontSize=32,
                    leading=38,
                    alignment=1, # Center
                    textColor=colors.HexColor(theme["primary"])
                )
                ty_sub_style = ParagraphStyle(
                    'TYSub',
                    parent=subtitle_style,
                    fontSize=14,
                    leading=18,
                    alignment=1, # Center
                    textColor=colors.HexColor(theme["accent"])
                )

                slide_elements.append(Spacer(1, 15))
                slide_elements.append(Paragraph("Thank You!", ty_h_style))
                slide_elements.append(Spacer(1, 6))
                slide_elements.append(Paragraph(clean_md_to_reportlab(strip_emojis_for_pdf(s.subtitle or "Questions & Classroom Discussion")), ty_sub_style))
                slide_elements.append(Spacer(1, 20))

                ty_bullets = []
                for b in s.bullets:
                    ty_bullets.append(Paragraph(f"• {clean_md_to_reportlab(strip_emojis_for_pdf(b))}", bullet_style))
                    ty_bullets.append(Spacer(1, 6))

                ty_bullets.append(Spacer(1, 10))
                ty_bullets.append(Paragraph(f"<b>Presented by:</b> {presenter_display}   •   <b>DEVGYA Global Edutech Private Limited</b>", cat_badge_style))

                card_table = Table([[ty_bullets]], colWidths=[760])
                card_table.setStyle(TableStyle([
                    ('BACKGROUND', (0,0), (-1,-1), colors.HexColor(theme["card_bg"])),
                    ('BOX', (0,0), (-1,-1), 1.5, colors.HexColor(theme["accent"])),
                    ('PADDING', (0,0), (-1,-1), 16),
                    ('ALIGN', (0,0), (-1,-1), 'CENTER'),
                ]))
                slide_elements.append(card_table)

                if s.speaker_notes:
                    slide_elements.append(Spacer(1, 12))
                    notes_clean = clean_md_to_reportlab(strip_emojis_for_pdf(s.speaker_notes))
                    notes_table = Table([[Paragraph(f"<b>TEACHER CLOSING NOTES:</b> {notes_clean}", notes_style)]], colWidths=[760])
                    notes_table.setStyle(TableStyle([
                        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor("#F1F5F9")),
                        ('BOX', (0,0), (-1,-1), 0.5, colors.HexColor("#94A3B8")),
                        ('PADDING', (0,0), (-1,-1), 6),
                    ]))
                    slide_elements.append(notes_table)

                story.append(KeepTogether(slide_elements))
                if idx < total_slides - 1:
                    from reportlab.platypus import PageBreak
                    story.append(PageBreak())
                continue

            # 3. MIDDLE SLIDES
            # Slide Top Bar Table
            top_bar_data = [
                [
                    Paragraph(f"<b>DEVGYA GLOBAL EDUTECH</b> • {strip_emojis_for_pdf(pres_data.topic).upper()}", cat_badge_style),
                    Paragraph(f"<b>Slide {s.slide_number} of {total_slides}</b>", cat_badge_style)
                ]
            ]
            top_table = Table(top_bar_data, colWidths=[550, 220])
            top_table.setStyle(TableStyle([
                ('ALIGN', (0,0), (0,0), 'LEFT'),
                ('ALIGN', (1,0), (1,0), 'RIGHT'),
                ('PADDING', (0,0), (-1,-1), 0),
                ('BOTTOMPADDING', (0,0), (-1,-1), 4),
            ]))
            slide_elements.append(top_table)
            slide_elements.append(HRFlowable(width="100%", thickness=2, color=colors.HexColor(theme["accent"]), spaceBefore=2, spaceAfter=8))

            # Category & Slide Title
            slide_elements.append(Paragraph(strip_emojis_for_pdf(s.category).upper(), cat_badge_style))
            slide_elements.append(Spacer(1, 2))
            clean_title = clean_md_to_reportlab(strip_emojis_for_pdf(s.title))
            slide_elements.append(Paragraph(clean_title, title_style))
            if s.subtitle:
                clean_sub = clean_md_to_reportlab(strip_emojis_for_pdf(s.subtitle))
                slide_elements.append(Paragraph(clean_sub, subtitle_style))
            slide_elements.append(Spacer(1, 10))

            # Slide Content Card
            content_rows = []
            if s.layout == "two_column" and (s.left_column or s.right_column):
                l_title = (s.left_column or {}).get("title", "Aspect A")
                r_title = (s.right_column or {}).get("title", "Aspect B")
                l_bullets = "<br/>• ".join([clean_md_to_reportlab(strip_emojis_for_pdf(b)) for b in (s.left_column or {}).get("bullets", [])])
                r_bullets = "<br/>• ".join([clean_md_to_reportlab(strip_emojis_for_pdf(b)) for b in (s.right_column or {}).get("bullets", [])])

                cell_l = Paragraph(f"<b>{clean_md_to_reportlab(l_title)}</b><br/><br/>• {l_bullets}", bullet_style)
                cell_r = Paragraph(f"<b>{clean_md_to_reportlab(r_title)}</b><br/><br/>• {r_bullets}", bullet_style)
                content_rows.append([cell_l, cell_r])
                card_table = Table(content_rows, colWidths=[380, 380])
                card_table.setStyle(TableStyle([
                    ('BACKGROUND', (0,0), (-1,-1), colors.HexColor(theme["card_bg"])),
                    ('BOX', (0,0), (-1,-1), 1, colors.HexColor("#CBD5E1")),
                    ('INNERGRID', (0,0), (-1,-1), 0.5, colors.HexColor("#E2E8F0")),
                    ('PADDING', (0,0), (-1,-1), 10),
                    ('VALIGN', (0,0), (-1,-1), 'TOP'),
                ]))
                slide_elements.append(card_table)

            elif s.layout == "stat_highlight" and s.metrics:
                metric_cells = []
                for m in s.metrics[:3]:
                    val = clean_md_to_reportlab(strip_emojis_for_pdf(str(m.get("value", ""))))
                    lbl = clean_md_to_reportlab(strip_emojis_for_pdf(str(m.get("label", ""))))
                    metric_cells.append(Paragraph(f"<font size=22 color='{theme['accent']}'><b>{val}</b></font><br/><br/><b>{lbl}</b>", bullet_style))
                while len(metric_cells) < 3:
                    metric_cells.append(Paragraph("", bullet_style))

                card_table = Table([metric_cells], colWidths=[250, 250, 260])
                card_table.setStyle(TableStyle([
                    ('BACKGROUND', (0,0), (-1,-1), colors.HexColor(theme["card_bg"])),
                    ('BOX', (0,0), (-1,-1), 1, colors.HexColor(theme["accent"])),
                    ('INNERGRID', (0,0), (-1,-1), 0.5, colors.HexColor("#E2E8F0")),
                    ('PADDING', (0,0), (-1,-1), 16),
                    ('ALIGN', (0,0), (-1,-1), 'CENTER'),
                ]))
                slide_elements.append(card_table)

            elif s.layout == "process_timeline" and s.timeline_steps:
                step_cells = []
                for i, st in enumerate(s.timeline_steps[:3]):
                    st_title = clean_md_to_reportlab(strip_emojis_for_pdf(str(st.get("title") or st.get("step") or f"Phase {i+1}")))
                    st_desc = clean_md_to_reportlab(strip_emojis_for_pdf(str(st.get("description") or "")))
                    step_cells.append(Paragraph(f"<font size=11 color='{theme['accent']}'><b>STEP {i+1}</b></font><br/><br/><b>{st_title}</b><br/><br/>{st_desc}", bullet_style))
                while len(step_cells) < 3:
                    step_cells.append(Paragraph("", bullet_style))

                card_table = Table([step_cells], colWidths=[250, 250, 260])
                card_table.setStyle(TableStyle([
                    ('BACKGROUND', (0,0), (-1,-1), colors.HexColor(theme["card_bg"])),
                    ('BOX', (0,0), (-1,-1), 1, colors.HexColor(theme["accent"])),
                    ('INNERGRID', (0,0), (-1,-1), 0.5, colors.HexColor("#E2E8F0")),
                    ('PADDING', (0,0), (-1,-1), 14),
                    ('VALIGN', (0,0), (-1,-1), 'TOP'),
                ]))
                slide_elements.append(card_table)

            else:
                # Main Bullet Card + Visual Callout
                bullets_p = []
                for b in s.bullets:
                    bullets_p.append(Paragraph(f"• {clean_md_to_reportlab(strip_emojis_for_pdf(b))}", bullet_style))
                    bullets_p.append(Spacer(1, 4))

                img_data = _download_image_bytes(s.image_url) if s.image_url else None
                pic_inserted = False
                if img_data:
                    try:
                        rl_pic = RLImage(io.BytesIO(img_data), width=230, height=135)
                        callout_p = [
                            rl_pic,
                            Spacer(1, 4),
                            Paragraph(clean_md_to_reportlab(strip_emojis_for_pdf(s.image_caption or s.image_keyword or "Visual Guide")), notes_style)
                        ]
                        pic_inserted = True
                    except Exception as img_err:
                        logger.warning(f"PDF slide image rendering failed: {img_err}")
                        pic_inserted = False

                if not pic_inserted:
                    callout_p = [
                        Paragraph(f"<b>VISUAL CONCEPT FOCUS:</b>", cat_badge_style),
                        Spacer(1, 4),
                        Paragraph(clean_md_to_reportlab(strip_emojis_for_pdf(s.image_caption or s.image_keyword)), bullet_style),
                        Spacer(1, 8),
                        Paragraph(f"<i>Keywords: {strip_emojis_for_pdf(s.image_keyword)}</i>", notes_style)
                    ]

                card_table = Table([[bullets_p, callout_p]], colWidths=[510, 250])
                card_table.setStyle(TableStyle([
                    ('BACKGROUND', (0,0), (0,0), colors.HexColor(theme["card_bg"])),
                    ('BACKGROUND', (1,0), (1,0), colors.HexColor("#F8FAFC")),
                    ('BOX', (0,0), (-1,-1), 1, colors.HexColor("#CBD5E1")),
                    ('INNERGRID', (0,0), (-1,-1), 1, colors.HexColor("#E2E8F0")),
                    ('PADDING', (0,0), (-1,-1), 10),
                    ('VALIGN', (0,0), (-1,-1), 'TOP'),
                ]))
                slide_elements.append(card_table)

            # Teacher Speaker Notes Box at bottom
            if s.speaker_notes:
                slide_elements.append(Spacer(1, 8))
                notes_clean = clean_md_to_reportlab(strip_emojis_for_pdf(s.speaker_notes))
                notes_table = Table(
                    [[Paragraph(f"<b>TEACHER PEDAGOGY GUIDANCE / SPEAKER NOTES:</b> {notes_clean}", notes_style)]],
                    colWidths=[760]
                )
                notes_table.setStyle(TableStyle([
                    ('BACKGROUND', (0,0), (-1,-1), colors.HexColor("#F1F5F9")),
                    ('BOX', (0,0), (-1,-1), 0.5, colors.HexColor("#94A3B8")),
                    ('PADDING', (0,0), (-1,-1), 6),
                ]))
                slide_elements.append(notes_table)

            story.append(KeepTogether(slide_elements))

            # Pagebreak between slides (except last)
            if idx < total_slides - 1:
                from reportlab.platypus import PageBreak
                story.append(PageBreak())

        doc.build(story)
        pdf_bytes = buffer.getvalue()
        buffer.close()
        return pdf_bytes

ppt_service = PPTGeneratorService()
